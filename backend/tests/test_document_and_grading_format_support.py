from __future__ import annotations

import io
from pathlib import Path
from types import SimpleNamespace

import pytest
from fastapi import HTTPException
from pptx import Presentation
from starlette.datastructures import UploadFile

from api.routes import assignments
from models.assignment import Submission, SubmissionStatus
from models.user import UserRole
from workers.embedding_task import parse_resource_content
from workers.grading_task import _build_grading_content


REPO_ROOT = Path(__file__).resolve().parents[2]


def _sample_payload(file_type: str) -> bytes:
    if file_type == "pdf":
        return (REPO_ROOT / "data" / "dataStructure" / "ds01.pdf").read_bytes()
    if file_type == "docx":
        return (
            REPO_ROOT
            / "data"
            / "计算机科学与技术专业导论（实验报告）"
            / "Excel实验报告-批改.docx"
        ).read_bytes()
    if file_type == "xlsx":
        return (
            REPO_ROOT
            / "data"
            / "计算机科学与技术专业导论（实验报告）"
            / "Excel实验原始素材文件.xlsx"
        ).read_bytes()
    if file_type == "pptx":
        stream = io.BytesIO()
        deck = Presentation()
        slide = deck.slides.add_slide(deck.slide_layouts[1])
        slide.shapes.title.text = "EduAI PPTX Parsing"
        slide.placeholders[1].text = "Format coverage evidence for grading input."
        deck.save(stream)
        return stream.getvalue()
    raise AssertionError(f"Unsupported test file type: {file_type}")


class _FakeMinioResponse:
    def __init__(self, payload: bytes):
        self._payload = payload

    def read(self):
        return self._payload

    def close(self):
        return None

    def release_conn(self):
        return None


class _FakeMinioClient:
    def __init__(self, payload: bytes):
        self._payload = payload

    def get_object(self, bucket, object_name):
        return _FakeMinioResponse(self._payload)


@pytest.mark.parametrize("file_type", ["pdf", "docx", "pptx"])
def test_parse_resource_content_extracts_text_from_rbs_formats(file_type: str):
    payload = _sample_payload(file_type)

    text = parse_resource_content(file_type, payload).strip()

    assert text
    assert len(text) >= 20


@pytest.mark.parametrize("file_type", ["pdf", "docx", "pptx", "xlsx"])
def test_build_grading_content_uses_supported_attachment_formats(monkeypatch, file_type: str):
    payload = _sample_payload(file_type)
    expected = parse_resource_content(file_type, payload).strip()
    submission = Submission(content=None, file_path=f"submissions/demo/sample.{file_type}")

    monkeypatch.setattr(
        "workers.embedding_task.get_minio_client",
        lambda: _FakeMinioClient(payload),
    )

    merged, warnings = _build_grading_content(submission)

    assert warnings == ""
    assert "[附件解析内容]" in merged
    assert expected[: min(len(expected), 200)] in merged


@pytest.mark.asyncio
async def test_get_grading_result_returns_worker_failure_prompt(monkeypatch):
    submission = Submission(status=SubmissionStatus.FAILED)

    async def fake_get_submission_context(db, submission_id):
        return submission, SimpleNamespace(), SimpleNamespace()

    monkeypatch.setattr(assignments, "_get_submission_context", fake_get_submission_context)
    monkeypatch.setattr(assignments, "ensure_submission_access", lambda submission, course, user: None)

    with pytest.raises(HTTPException) as exc:
        await assignments.get_grading_result(
            submission_id=1,
            db=SimpleNamespace(),
            user=SimpleNamespace(),
        )

    assert exc.value.status_code == 409
    assert exc.value.detail == "Grading failed"


@pytest.mark.asyncio
async def test_submit_assignment_returns_clear_upload_failure_prompt(monkeypatch):
    async def fake_get_assignment_course(db, assignment_id):
        return SimpleNamespace(id=assignment_id, course_id=3), SimpleNamespace(id=3, teacher_id=7)

    async def fake_ensure_course_access(db, *, course, user):
        return None

    async def fake_capability(db, *, course_id: int):
        return SimpleNamespace(has_grading=True)

    def fake_upload_bytes(*, object_name, data, content_type):
        raise RuntimeError("minio unavailable")

    monkeypatch.setattr(assignments, "_get_course_for_assignment", fake_get_assignment_course)
    monkeypatch.setattr(assignments, "ensure_course_access", fake_ensure_course_access)
    monkeypatch.setattr(assignments, "get_published_course_agent_capability", fake_capability)
    monkeypatch.setattr(assignments, "upload_bytes", fake_upload_bytes)

    file = UploadFile(filename="answer.pdf", file=io.BytesIO(b"fake payload"))

    with pytest.raises(HTTPException) as exc:
        await assignments.submit_assignment(
            assignment_id=5,
            content=None,
            file=file,
            db=SimpleNamespace(),
            user=SimpleNamespace(id=42, role=UserRole.STUDENT),
        )

    assert exc.value.status_code == 500
    assert exc.value.detail == "Failed to upload submission file: minio unavailable"


@pytest.mark.asyncio
async def test_submit_assignment_rejects_unsupported_file_type(monkeypatch):
    async def fake_get_assignment_course(db, assignment_id):
        return SimpleNamespace(id=assignment_id, course_id=3), SimpleNamespace(id=3, teacher_id=7)

    async def fake_ensure_course_access(db, *, course, user):
        return None

    async def fake_capability(db, *, course_id: int):
        return SimpleNamespace(has_grading=True)

    monkeypatch.setattr(assignments, "_get_course_for_assignment", fake_get_assignment_course)
    monkeypatch.setattr(assignments, "ensure_course_access", fake_ensure_course_access)
    monkeypatch.setattr(assignments, "get_published_course_agent_capability", fake_capability)

    file = UploadFile(filename="answer.zip", file=io.BytesIO(b"fake payload"))

    with pytest.raises(HTTPException) as exc:
        await assignments.submit_assignment(
            assignment_id=5,
            content=None,
            file=file,
            db=SimpleNamespace(),
            user=SimpleNamespace(id=42, role=UserRole.STUDENT),
        )

    assert exc.value.status_code == 400
    assert exc.value.detail == "不支持该文件类型：.zip"
