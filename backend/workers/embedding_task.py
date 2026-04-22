from __future__ import annotations

import base64
import re
from io import BytesIO

import dashscope
from minio import Minio
from docx import Document
from openpyxl import load_workbook
from PyPDF2 import PdfReader
from pptx import Presentation
from sqlalchemy import create_engine, delete, select
from sqlalchemy.orm import sessionmaker
import pypdfium2 as pdfium

from core.config import get_settings
from models.course import CourseResource, ResourceChunk
from workers.celery_app import celery_app

settings = get_settings()
sync_engine = create_engine(settings.DATABASE_SYNC_URL)
SyncSessionLocal = sessionmaker(bind=sync_engine)


def get_minio_client() -> Minio:
    return Minio(
        settings.MINIO_ENDPOINT,
        access_key=settings.MINIO_ACCESS_KEY,
        secret_key=settings.MINIO_SECRET_KEY,
        secure=False,
    )


def split_text(content: str, chunk_size: int = 500, overlap: int = 100) -> list[str]:
    text = normalize_text(content)
    if not text:
        return []

    paragraphs = [part.strip() for part in re.split(r"\n\s*\n+", text) if part.strip()]
    if not paragraphs:
        paragraphs = [text]

    chunks: list[str] = []
    current = ""

    for paragraph in paragraphs:
        if len(paragraph) > chunk_size:
            if current:
                chunks.append(current)
                current = ""
            chunks.extend(split_large_block(paragraph, chunk_size, overlap))
            continue

        candidate = f"{current}\n\n{paragraph}".strip() if current else paragraph
        if len(candidate) <= chunk_size:
            current = candidate
        else:
            if current:
                chunks.append(current)
            current = paragraph

    if current:
        chunks.append(current)

    return chunks


def parse_resource_content(file_type: str, payload: bytes) -> str:
    suffix = file_type.lower()
    if suffix in {"txt", "md", "py", "json", "csv"}:
        return payload.decode("utf-8", errors="ignore")

    if suffix == "pdf":
        return parse_pdf_content(payload)

    if suffix == "docx":
        return parse_docx_content(payload)

    if suffix == "pptx":
        return parse_pptx_content(payload)

    if suffix == "xlsx":
        return parse_xlsx_content(payload)

    return payload.decode("utf-8", errors="ignore")


def parse_xlsx_content(payload: bytes) -> str:
    workbook = load_workbook(BytesIO(payload), data_only=True)
    texts: list[str] = []

    try:
        for worksheet in workbook.worksheets:
            rows: list[str] = []
            for row in worksheet.iter_rows(values_only=True):
                values = [str(value).strip() for value in row if value not in (None, "")]
                if values:
                    rows.append(" | ".join(values))

            if rows:
                texts.append(f"[{worksheet.title}]\n" + "\n".join(rows))
    finally:
        workbook.close()

    return "\n\n".join(texts)


def parse_docx_content(payload: bytes) -> str:
    document = Document(BytesIO(payload))
    texts: list[str] = []

    texts.extend(paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip())

    for table in document.tables:
        for row in table.rows:
            values = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if values:
                texts.append(" | ".join(values))

    return "\n\n".join(texts)


def parse_pptx_content(payload: bytes) -> str:
    presentation = Presentation(BytesIO(payload))
    texts: list[str] = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        slide_texts: list[str] = []
        for shape in slide.shapes:
            slide_texts.extend(extract_shape_text(shape))

        if slide_texts:
            texts.append(f"[Slide {slide_index}]\n" + "\n".join(slide_texts))

    return "\n\n".join(texts)


def parse_pdf_content(payload: bytes) -> str:
    reader = PdfReader(BytesIO(payload))
    texts = [(page.extract_text() or "").strip() for page in reader.pages]
    extracted_text = "\n".join(text for text in texts if text)
    if is_meaningful_text(extracted_text):
        return extracted_text

    return ocr_pdf_content(payload)


def ocr_pdf_content(payload: bytes) -> str:
    if not settings.DASHSCOPE_API_KEY:
        return ""

    dashscope.api_key = settings.DASHSCOPE_API_KEY
    pdf = pdfium.PdfDocument(BytesIO(payload))
    texts: list[str] = []

    try:
        page_total = len(pdf)
        max_pages = settings.PDF_OCR_MAX_PAGES
        page_count = page_total if max_pages <= 0 else min(page_total, max_pages)

        for page_index in range(page_count):
            try:
                page = pdf[page_index]
                pil_image = page.render(scale=2).to_pil()
                buffer = BytesIO()
                pil_image.save(buffer, format="PNG")
                image_base64 = base64.b64encode(buffer.getvalue()).decode("utf-8")

                response = dashscope.MultiModalConversation.call(
                    model=settings.QWEN_VL_MODEL,
                    messages=[
                        {
                            "role": "user",
                            "content": [
                                {"image": f"data:image/png;base64,{image_base64}"},
                                {"text": "请提取这页中的文字内容，直接返回纯文本，不要解释。"},
                            ],
                        }
                    ],
                )
            except Exception as exc:
                print(f"[Embedding][OCR] page={page_index + 1} failed before response: {exc}")
                continue

            if response.status_code != 200:
                print(
                    f"[Embedding][OCR] page={page_index + 1} request failed: "
                    f"{getattr(response, 'code', None)} {getattr(response, 'message', None)}"
                )
                continue

            try:
                content = response.output.choices[0].message.content
                page_text_parts: list[str] = []
                if isinstance(content, list):
                    for item in content:
                        if isinstance(item, dict) and item.get("text"):
                            page_text_parts.append(item["text"].strip())
                elif isinstance(content, str):
                    page_text_parts.append(content.strip())
            except Exception as exc:
                print(f"[Embedding][OCR] page={page_index + 1} response parse failed: {exc}")
                continue

            page_text = "\n".join(part for part in page_text_parts if part)
            if page_text:
                texts.append(page_text)
    finally:
        pdf.close()

    return "\n\n".join(texts)


@celery_app.task(name="workers.embedding_task.process_resource")
def process_resource(resource_id: int):
    """处理课程资料并将其分块写入 resource_chunks。"""
    print(f"[Embedding] 开始处理资料 resource_id={resource_id}")

    with SyncSessionLocal() as db:
        resource = db.execute(
            select(CourseResource).where(CourseResource.id == resource_id)
        ).scalar_one_or_none()
        if not resource:
            return {"resource_id": resource_id, "status": "not_found"}

        try:
            resource.processing_status = "processing"
            resource.processing_error = None
            resource.is_processed = False
            db.commit()

            client = get_minio_client()
            response = client.get_object(settings.MINIO_BUCKET, resource.file_path)
            try:
                payload = response.read()
            finally:
                response.close()
                response.release_conn()

            content = parse_resource_content(resource.file_type, payload)
            chunks = split_text(content)

            db.execute(delete(ResourceChunk).where(ResourceChunk.resource_id == resource.id))

            for index, chunk in enumerate(chunks):
                db.add(
                    ResourceChunk(
                        resource_id=resource.id,
                        course_id=resource.course_id,
                        content=chunk,
                        chunk_index=index,
                        metadata_={
                            "resource_name": resource.name,
                            "file_type": resource.file_type,
                            "object_name": resource.file_path,
                        },
                    )
                )

            resource.chunk_count = len(chunks)
            resource.is_processed = True
            resource.processing_status = "processed"
            resource.processing_error = None
            db.commit()
            return {"resource_id": resource_id, "status": "processed", "chunk_count": len(chunks)}
        except Exception as exc:
            db.rollback()
            resource = db.execute(
                select(CourseResource).where(CourseResource.id == resource_id)
            ).scalar_one_or_none()
            if resource:
                resource.chunk_count = 0
                resource.is_processed = False
                resource.processing_status = "failed"
                resource.processing_error = str(exc)
                db.commit()

            print(f"[Embedding] 处理失败 resource_id={resource_id}: {exc}")
            return {"resource_id": resource_id, "status": "failed", "error": str(exc)}


def normalize_text(content: str) -> str:
    normalized = content.replace("\r\n", "\n").replace("\r", "\n")
    normalized = re.sub(r"[ \t]+\n", "\n", normalized)
    normalized = re.sub(r"\n{3,}", "\n\n", normalized)
    return normalized.strip()


def split_large_block(text: str, chunk_size: int, overlap: int) -> list[str]:
    chunks: list[str] = []
    step = max(chunk_size - overlap, 1)
    start = 0

    while start < len(text):
        end = min(start + chunk_size, len(text))
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        if end >= len(text):
            break
        start += step

    return chunks


def is_meaningful_text(text: str) -> bool:
    normalized = normalize_text(text)
    if len(normalized) < settings.PDF_TEXT_MIN_LENGTH:
        return False

    meaningful_chars = sum(char.isalnum() or "\u4e00" <= char <= "\u9fff" for char in normalized)
    ratio = meaningful_chars / max(len(normalized), 1)
    return ratio >= settings.PDF_TEXT_MIN_MEANINGFUL_RATIO


def extract_shape_text(shape) -> list[str]:
    texts: list[str] = []

    if hasattr(shape, "text") and shape.text:
        text = shape.text.strip()
        if text:
            texts.append(text)

    if hasattr(shape, "shapes"):
        for child in shape.shapes:
            texts.extend(extract_shape_text(child))

    return texts
